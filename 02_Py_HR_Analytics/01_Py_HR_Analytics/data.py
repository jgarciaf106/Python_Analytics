# library imports
import os
import codecs
import pyodbc
import json
import argparse
import re
import win32com.client
import datetime as dt
import pandas as pd
import numpy as np
import xlwings as xw
import tkinter as tk
import matplotlib.pyplot as plt
import seaborn as sns
import datetime as dt
import fiscalyear as fy
import config as cfg
import utils as ut
from sys import exit
from tkinter import messagebox
from sqlalchemy import create_engine
from openpyxl import Workbook
from openpyxl.utils.dataframe import dataframe_to_rows
from openpyxl.styles import PatternFill, Border, Side, Alignment, Protection, Font
from win32com.client import Dispatch, constants
from pptx import Presentation
from pptx.util import Inches


class Data:

    # init method or constructor
    def __init__(self,):
        self.__db_type = cfg.db["type"]
        self.__db_name = cfg.db["database"]
        self.__output_path = None
        self.__load_path = None
        self.__query_path = None
        self.__input_path = None
        self.__today = dt.datetime.today() - pd.DateOffset(months=1)
        self.__monthname_long = self.__today.strftime("%B")
        self.__monthname_short = self.__today.strftime("%b")
        self.__monthyear = self.__today.strftime("%m%y")
        self.__fiscal_year = None
        self.__current_quarter = None
        self.__fiscal_quarter = None
        self.__engine = self.__set_engine()
        self.__set_fiscal_year()

    # setter functions

    def set_path(self, custom=False, path_type=None, custom_path=None):
        """
        Sets paths to specific locations, input, output, queries

        Parameters
        ----------
        path_type : str
            Whether the path is for an input, output or query location
        custom_path : str
            Path to the specific directory location
        """

        if not custom:
            self.__input_path = cfg.path["input"]
            self.__load_path = cfg.path["load"]
            self.__output_path = cfg.path["output"]
            self.__query_path = cfg.path["queries"]
        else:
            if path_type == "output":
                self.__output_path = custom_path
            elif path_type == "load":
                self.__load_path = custom_path
            elif path_type == "query":
                self.__query_path = custom_path
            elif path_type == "input":
                self.__input_path = custom_path

    def __set_fiscal_year(self, f_year=None, s_month=11, s_day=1, s_year="previous"):
        """
        Sets the fiscal year to an company's Fiscal start date.

        Parameters
        ----------
        f_year : int, optional
            Year when fiscal year starts (default is
            blank)
        s_month : int, optional
            Month when fiscal year starts (default is
            11)
        s_day : int, optional
            Day when fiscal year starts (default is
            01)
        s_year : str, optional
            Point in time when fiscal year starts, same year of previous year (default is
            previous)
        """

        if f_year is None:
            if dt.datetime.today().month < 11:
                f_year = dt.datetime.today().year - 1
            else:
                f_year = dt.datetime.today().year

        self.__fiscal_year = fy.FiscalYear(f_year)
        fy.setup_fiscal_calendar(
            start_year=s_year, start_month=s_month, start_day=s_day
        )
        self.__fiscal_quarter = str(
            fy.FiscalQuarter.current().prev_fiscal_quarter
        ).split(" ")[0]
        self.__current_quarter = str(
            fy.FiscalQuarter.current().prev_fiscal_quarter
        ).split(" ")[1]

    # set engine db
    def __set_engine(self):
        """
        Sets the engine to start querying the specific data base

        Parameters
        ----------
        None

        Returns
        -------
        object
            details of the database that needs to be query.
        """

        # db create connection engine
        if self.__db_type == "postgres":
            return create_engine(
                "postgresql://{0}:{1}@localhost/{2}".format(
                    cfg.db["user"], cfg.db["password"], self.__db_name
                )
            )
        elif self.__db_type == "mssql":
            return create_engine(
                f"mssql+pyodbc://HP-308027A\\BIHRDW/{self.__db_name}?trusted_connection=yes&driver=ODBC+Driver+17+for+SQL+Server"
            )
        elif self.__db_type == "oracle":
            return create_engine(
                "'oracle://{0}:{1}@localhost/{2}".format(
                    cfg.db["user"], cfg.db["password"], self.__db_name
                )
            )
        elif self.__db_type == "mysql":
            return create_engine(
                "mysql://{0}:{1}@localhost/{2}".format(
                    cfg.db["user"], cfg.db["password"], self.__db_name
                )
            )

    # query external file queries
    def external_query(self, query_name):
        """
        Query the database using a external query file.

        Parameters
        ----------
        query_name : str
            The file with the predefined query.

        Returns
        -------
        dataframe
            a dataframe with the queried data.
        """

        query_path = self.__query_path + "{0}.sql".format(query_name)
        with open(query_path) as get:
            query = get.read()
        return pd.read_sql_query(query, self.__engine)

    # query inline queries
    def internal_query(self, query_name):
        """
        Query the database using an inline query or query variable.

        Parameters
        ----------
        query_name : str
            The file with the predefined query.

        Returns
        -------
        dataframe
            a dataframe with the queried data.
        """

        return pd.read_sql_query(query_name, self.__engine)

    def execute_query(self, query_name):
        with self.__engine.connect() as connection:
            connection.execute(query_name)

    # read  files
    def file_query(self, file_name, engine_reader=None, read_sheet=0):
        """
        Reads an external excel file to get the data

        Parameters
        ----------
        file_name : str
            The file name that needs to be read.
        engine_reader : str, optional
            If a a different engine reader is neeed to open the file (default is
            Blank)
        read_sheet : [str, int], optional
            Indicates which sheet to read by Sheet name or Index position (default is
            index 0)

        Returns
        -------
        dataframe
            a dataframe with the queried data.
        """

        file = self.__input_path + file_name
        if engine_reader is None:
            return pd.read_excel(file, sheet_name=read_sheet, index_col=0)
        else:
            return pd.read_excel(
                file, sheet_name=read_sheet, index_col=0, engine=engine_reader
            )

    def __password_tracker(self, request_date, file_name, password):
        """
        Stores a password for all password protected files in a tracker

        Parameters
        ----------
        request_date : date
            The date when the request was process
        file_name : string
            The request file's name
        password : string
            The password to open the file
        """

        app = xw.App(visible=False)
        wb = xw.Book(
            r"C:/Users/garciand/OneDrive - HP Inc/Desktop/Deliverables/Output/Password Tracker/Password_Tracker.xlsx"
        )
        ws = wb.sheets[0]

        next_row = "A" + str(ws.range("A1").current_region.last_cell.row + 1)

        ws.range(next_row).options(index=False).value = [
            request_date,
            file_name,
            self.__output_path,
            password,
        ]

        wb.save()
        app.quit()

    # password share after email sent
    def __password_share(self, password, email_subject):
        """
        Replies to a sent email with the file password.

        Parameters
        ----------
        password : str
            The file's password
        email_subject : str
            Subject of the last email to ensure the password is sent only to that email.

        Returns
        -------
        list
            a list of strings used that are the header columns
        """

        outlook = win32com.client.Dispatch("Outlook.Application").GetNamespace("MAPI")
        sent_items = outlook.GetDefaultFolder(5)
        messages = sent_items.Items
        message = messages.GetLast()
        pass

    # save email to outlook drafts
    def __draft_email(self, recipient, file_name, password, folder_search):
        """
        Creates a draft email to be sent with the report atteched.

        Parameters
        ----------
        recipient : str
            Email address of the person that should get the email
        file_name : str
            The file location and name to be attached
        password : str
            The password to open the file
        folder_search : str, optional
            Determines in which email folder the last email is stores (default is
            Blank)

        Returns
        -------
        None
            define a null variable or an object
        """

        # instantiate outlook to get received Items
        outlook = win32com.client.Dispatch("Outlook.Application").GetNamespace("MAPI")
        if folder_search is not None:
            inbox = outlook.GetDefaultFolder(6).Folders.Item(
                folder_search
            )  # inbox in subfolders
        else:
            inbox = outlook.GetDefaultFolder(6)  # inbox emails

        # instantiate outlook to send email
        const = win32com.client.constants
        olMailItem = 0x0
        obj = win32com.client.Dispatch("Outlook.Application")

        # get recived emails
        messages = inbox.Items
        message = messages.GetLast()
        sender_name = message.Sender.GetExchangeUser().name.split(",")[-1].strip()
        sender_address = message.Sender.GetExchangeUser().PrimarySmtpAddress
        sender_cc_address = ";".join(
            [
                item.AddressEntry.GetExchangeUser().PrimarySmtpAddress
                for item in message.Recipients
                if item.AddressEntry.GetExchangeUser().PrimarySmtpAddress
                != "andres.garcia.fernandez@hp.com"
            ]
        )

        # set email signature
        signature_htm = os.path.join(
            os.environ["USERPROFILE"], "AppData\\Roaming\\Microsoft\\Signatures\\HP.htm"
        )
        html_file = codecs.open(signature_htm, "r", "utf-8", errors="ignore")
        email_signature = html_file.read()
        html_file.close()

        # html email body and subject

        # get original email subject
        email_subject = message.Subject

        # create new email items
        new_mail = obj.CreateItem(olMailItem)
        new_mail.Subject = email_subject
        new_mail.BodyFormat = 2

        if password is not None:
            email_body = """
                    <HTML>
                        <BODY style="font-family:HP Simplified Light;font-size:14.5px;">
                            <p>Hi {0},</p>
                            <p>Please find attached the report requested.</p>
                            <p>The attached file is password protected, password will be shared on the following email.</p>
                            <p><em style="color:red">*** Delete Before Sending Email ***</em>
                            <strong>Password:</strong> {1} Copy password to send it on a separate email.
                            <em style="color:red">*** Delete Before Sending Email ***</em></p>
                            <p>Kind Regards,</p> 
                        </BODY>
                    </HTML>""".format(
                sender_name, password
            )
        else:
            email_body = """
                <HTML>
                    <BODY style="font-family:HP Simplified Light;font-size:14.5px;">
                        <p>Hi {0},</p>
                        <p>Please find attached the report requested.</p>
                        <p>Kind Regards,</p> 
                    </BODY>
                </HTML>""".format(
                sender_name
            )

        # validate that email will be sent to correct recepient
        if sender_address == recipient:  # based on the subject replying to email
            reply_all = message.ReplyAll()
            new_mail.HTMLBody = email_body + reply_all.HTMLBody
            new_mail.To = sender_address
            new_mail.CC = sender_cc_address
            attachment = self.__output_path + file_name + ".xlsx"
            new_mail.Attachments.Add(Source=attachment)
            new_mail.save()

        return None

    # save formatted file
    def __file_saver(
        self, data, file_name, protect_file, draft_email, requester, email_folder
    ):
        """
        Saves a formatted excel file

        Parameters
        ----------
        data : dataframe
            The datafram that needs to be stored to excel
        file_name : str
            The file name for the excel file
        protect_file : bool
            Flag to determine if the file has to be password protected
        draft_email : bool
            Flag to determine if an email draft is needed
        requester : str
            email address if an email draft has to sent
        email_folder : str
            Email folder where the last email is stored
        """

        # default password
        file_password = None

        # remove existing files before save
        ut.file_cleaner("{0}{1}.xlsx".format(self.__output_path, file_name))

        # workbook / sheet variables
        app = xw.App(visible=False)
        wb = xw.Book()
        ws = wb.sheets[0]

        # excel data header formatting
        ws.range("A1").options(pd.DataFrame, index=False).value = ut.number_to_string(
            data, "Worker ID"
        )
        header_format = ws.range("A1").expand("right")
        header_format.color = cfg.file_format["header_bg_color"]
        header_format.api.Font.Name = cfg.file_format["font_name"]
        header_format.api.Font.Color = cfg.file_format["header_font_color"]
        header_format.api.Font.Bold = True
        header_format.api.Font.Size = cfg.file_format["header_font_size"]

        # excel data content formatting
        data_format = ws.range("A2").current_region
        data_format.api.Font.Name = cfg.file_format["font_name"]
        data_format.api.Font.Size = cfg.file_format["content_font_size"]

        # save password protect file if needed
        if protect_file:
            file_password = "HPI" + str(self.__today.strftime("%I%M%S"))
            wb.api.SaveAs(self.__output_path + file_name, Password=file_password)

        else:
            wb.api.SaveAs(self.__output_path + file_name)

        app.quit()

        # update password tracker
        self.__password_tracker(
            self.__today, file_name, file_password,
        )

        # prompt
        if draft_email:
            self.__draft_email(requester, file_name, file_password, email_folder)

    # export file to folder
    def export_data(
        self,
        odf,
        custom_filename=None,
        protect_file=False,
        draft_email=False,
        requester=None,
        email_folder=None,
    ):
        """
        Calls other functions to export the data.

        Parameters
        ----------
        odf : dataframe
            The dataframe that needs to be stored.
        custom_filename : str, optional
            A custome file name for the file to be stored (default is
            Blank)
        protect_file : bool, optional
            A flag used to determine if file has to password protected(default is
            False)
        draft_email : bool, optional
            A flag used to determine if a draft email is neeed (default is
            False)
        requester : str, optional
            Email address for the requester (default is
            Blank)
        email_folder : str, optional
            Email folder where the last email is stored (default is
            Blank)

        Returns
        -------
        list
            a list of strings used that are the header columns
        """

        # assign file name
        if custom_filename is None:
            file_name = "Output Report " + self.__today.strftime("%Y-%m-%d")
        else:
            file_name = custom_filename

        # save file
        self.__file_saver(
            odf, file_name, protect_file, draft_email, requester, email_folder
        )

    def ppt_export(self, file_type, template_type, org=None):
        """
        Exports the formatted ppt

        Parameters
        ----------
        file_type : str
            Determines if the ppt if HPI or L1 org
        template_type : str
            Determines which template to use dark or light
        org : str, optional
            A flag used to indicate the L1 Org Name (default is
            None)

        Returns
        -------
        list
            a list of strings used that are the header columns
        """

        # determine output file HPI Total or L1 Org Dashboard
        if file_type == "HPI":
            output_file = "HPI DEI Dashboard {0} {1}".format(
                self.__current_quarter, self.__fiscal_quarter
            )
        elif file_type == "L1 ORG":
            output_file = "DEI {0} Dashboard {1}".format(org, self.__monthyear)

        # set slides to be updated
        prs = Presentation(
            "../Templates/HP_Presentation_Template_" + template_type + ".pptx"
        )
        slide_1 = prs.slides[0]
        slide_4 = prs.slides[3]
        slide_5 = prs.slides[4]
        slide_6 = prs.slides[5]
        slide_7 = prs.slides[6]
        slide_8 = prs.slides[7]
        slide_9 = prs.slides[8]
        slide_11 = prs.slides[10]
        slide_12 = prs.slides[11]
        slide_13 = prs.slides[12]
        slide_15 = prs.slides[14]

        # set placeholders to be updated
        prs_sub_1 = slide_1.placeholders[1]
        prs_sub_4 = slide_4.placeholders[0]
        prs_sub_6 = slide_6.placeholders[0]
        prs_sub_7 = slide_7.placeholders[0]
        prs_sub_8 = slide_8.placeholders[0]
        prs_sub_9 = slide_9.placeholders[0]

        # update placeholders
        if file_type == "HPI":
            prs_sub_1.text = "As of {0} month, end/{1}".format(
                self.__monthname_long, self.__current_quarter
            )
        elif file_type == "L1 ORG":
            prs_sub_1.text = "{0} (As of {1} month, end/{2})".format(
                org, self.__monthname_long, self.__current_quarter
            )

        prs_sub_4.text = (
            self.__fiscal_quarter + " " + self.__current_quarter + " Headcount"
        )

        prs_sub_6.text = "{0}/{1} Status to Diversity Targets (Company Level)".format(
            self.__monthname_short, self.__current_quarter
        )
        prs_sub_7.text = "{0}/{1} Active Headcount by Organization / MRU".format(
            self.__monthname_short, self.__current_quarter
        )
        prs_sub_8.text = "{0}/{1} Active Headcount by Organization / MRU (Absolute values)".format(
            self.__monthname_short, self.__current_quarter
        )
        prs_sub_9.text = "{0}/{1} US Ethnic Groups Not Self-Identified HC by Organization / MRU".format(
            self.__monthname_short, self.__current_quarter
        )

        # save updated template
        prs.save(
            self.__output_path
            + "Quarterly Dashboards\\"
            + file_type
            + "\\"
            + output_file
            + ".pptx"
        )

