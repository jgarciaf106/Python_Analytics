import os
from pptx import Presentation

# delete existing file
def file_cleaner(file_name):
    """
    Deletes any old file before storing an updated file

    Parameters
    ----------
    file_name : str
        The file location and name to be deleted
    """

    if os.path.exists(file_name):
        os.remove(file_name)


# trims whitespaces
def column_trim(df):
    """
    Trims any trailing spaces on string dataframe colums

    Parameters
    ----------
    df : dataframe
        The dataframe that needs column tims

    Returns
    -------
    datafram
        a dataframe with all column strings trimmed
    """

    trim_strings = lambda x: x.strip() if isinstance(x, str) else x
    return df.applymap(trim_strings)


def number_to_string(self, df, column_name):
    """
    Converts a column of numbers to strings

    Parameters
    ----------
    data : pandas.DataFrame
        The dataframe to be converted
    column_name : str
        The column name to be converted

    Returns
    -------
    pandas.DataFrame
        The converted dataframe
    """
    if column_name in df.columns:
        df[column_name] = "'" + df[column_name]
        return df
    else:
        return df


def ppt_identifier(self, input, slide_number):
    """
    Identifies a ppt single slide elements

    Parameters
    ----------
    input : str
        The ppt template directory location and name
    slide_number : int
        The slide number
    """

    prs = Presentation(input)
    slide = prs.slides[slide_number - 1]
    for shape in slide.shapes:
        print(
            "id: %s, name: %s, , type: %s"
            % (shape.shape_id, shape.name, shape.shape_type)
        )


def ppt_analyzer(self, input_path, output_path):
    """
    Analize a power point presentation to indicatate each of ppt elements

    Parameters
    ----------
    input : str
        The ppt template directory location and name
    output : str
        The file output to avoid overwrite the original template
    """

    prs = Presentation(input_path)

    for index, _ in enumerate(prs.slide_layouts):
        slide = prs.slides.add_slide(prs.slide_layouts[index])
        # Not every slide has to have a title
        try:
            title = slide.shapes.title
            title.text = "Title for Layout {}".format(index)
        except AttributeError:
            print("No Title for Layout {}".format(index))

        # Go through all the placeholders and identify them by index and type
        for shape in slide.placeholders:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                # Do not overwrite the title which is just a special placeholder
                try:
                    if "Title" not in shape.text:
                        shape.text = "Placeholder index:{} type:{}".format(
                            phf.idx, shape.name
                        )
                except AttributeError:
                    print("{} has no text attribute".format(phf.type))
                print("{} {}".format(phf.idx, shape.name))
    prs.save(output_path)
